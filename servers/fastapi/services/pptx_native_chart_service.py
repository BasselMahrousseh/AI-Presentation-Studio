import logging
import uuid
from typing import Any, Optional

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Length, Pt
from sqlmodel import select

from models.sql.slide import SlideModel
from services import chart_capture_store
from services.database import async_session_maker

LOGGER = logging.getLogger(__name__)

# The frontend capture reports chart geometry in the same 1280x720 logical
# canvas space every slide is designed and exported in (verified against a
# real pptx_model.json debug artifact: picture shape positions sit inside
# 0-1280 / 0-720 with no extra scale factor).
_DESIGN_WIDTH_PX = 1280
_DESIGN_HEIGHT_PX = 720

# A correct geometric match should be near-perfect (no scale correction is
# needed, see above), so a borderline score more likely means a bug (stale
# capture, wrong slide, cropped element) than a legitimate close call -> fail
# closed rather than tune a permissive threshold.
_IOU_MIN = 0.90
_IOU_AMBIGUITY_MARGIN = 0.05

Rect = tuple[int, int, int, int]  # left, top, width, height, all in EMU

# Chart kinds with no faithful native PPTX equivalent (polar_area, scatter,
# bubble) are intentionally absent here -> _resolve_xl_chart_type returns
# None for them and the caller leaves the flattened image untouched.
_FIXED_XL_CHART_TYPE_BY_KIND: dict[str, XL_CHART_TYPE] = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "horizontal_stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
}

_AXIS_TITLE_SUPPORTED_KINDS = {
    "bar",
    "horizontal_bar",
    "stacked_bar",
    "horizontal_stacked_bar",
    "line",
    "area",
}

# pie/donut are the only supported kinds with no category/value axis at all
# (chart.category_axis/.value_axis raise ValueError for them, confirmed
# against a real Presentation for every supported XL_CHART_TYPE).
_NO_AXIS_KINDS = {"pie", "donut"}

# Chart.js semantics: an array backgroundColor/borderColor on a dataset means
# one color per data point within that series (e.g. highlighting a single
# "leader" bar/slice a different color from its siblings). That reads
# naturally for these discrete/categorical kinds; line/area/radar keep a
# single color per series instead (a per-point-colored line reads as broken,
# not highlighted).
_POINT_COLOR_KINDS = {
    "bar",
    "horizontal_bar",
    "stacked_bar",
    "horizontal_stacked_bar",
    "pie",
    "donut",
}

# PowerPoint's default category-axis order for a horizontal bar chart lists
# the first category at the bottom; Chart.js lists it at the top. Reverse the
# axis so the native chart's category order visually matches the original.
_REVERSE_CATEGORY_ORDER_KINDS = {"horizontal_bar", "horizontal_stacked_bar"}

# PowerPoint rejects a data-label position on a doughnut chart outright - the
# file fails to open ("needs repair") and PowerPoint's own recovery drops the
# entire slide, confirmed by inspecting a real corrupted export: the chart8.xml
# doughnut part contained <c:dLblPos val="outEnd"/>, the sibling pie chart in
# the same file carried the identical value and opened fine, and a structural
# diff of the two parts found no other meaningful difference. python-pptx's
# own DataLabels.position setter (pptx/chart/datalabel.py) does zero
# chart-type validation, so nothing catches this before the file reaches
# PowerPoint - the gate below exists entirely to compensate for that. This is
# deliberately a set of one rather than growing to match _NO_AXIS_KINDS or
# _POINT_COLOR_KINDS: any future kind added here must be a conscious decision,
# not an accidental inclusion via a broader existing constant.
_OUTSIDE_END_LABEL_KINDS = {"pie"}

# DO NOT set data_labels.position for "area", even to XL_LABEL_POSITION.CENTER.
# Tried exactly that (CENTER is the one value Microsoft's own docs claim is
# valid on every chart type with no exception) to fix an area chart's labels
# rendering at the wrong height - and it corrupted a real exported file the
# same way the doughnut/OUTSIDE_END bug did: PowerPoint reported it needed
# repair, and repair deleted the area chart's own slide. So "ctr" is
# apparently *also* illegal for area in real PowerPoint, contradicting the
# documentation - confirmed only by a real user opening the real file, which
# is the one verification python-pptx itself can never provide (it accepted
# and re-read the file with zero complaint both times). Lesson: for this
# specific `dLblPos` element, trust nothing but an actual PowerPoint open -
# not vendor documentation, not python-pptx succeeding, not "this other
# value worked for a similar-sounding chart type". An area chart's data
# labels must be left entirely unpositioned (PowerPoint's own unconfigured
# default), even though that default visibly mis-places them - see the
# session notes in CLAUDE.md for what was tried and ruled out here before
# attempting anything further.

_XL_LEGEND_POSITION_BY_NAME: dict[str, XL_LEGEND_POSITION] = {
    "top": XL_LEGEND_POSITION.TOP,
    "left": XL_LEGEND_POSITION.LEFT,
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "right": XL_LEGEND_POSITION.RIGHT,
}

# The capture reports font sizes in CSS px within the same 1280x720 logical
# canvas the geometry above is measured in. That canvas is exported at
# 1280px = 13.333in (960pt), so 1 CSS px = 0.75pt exactly - and since every
# captured size is multiplied by a clean 0.75, the conversion is lossless in
# the centipoint units python-pptx's Font.size setter stores (e.g. 14px ->
# 10.5pt -> sz="1050").
_PT_PER_CSS_PX = 0.75
_MIN_FONT_PX = 1.0
_MAX_FONT_PX = 200.0


def _resolve_xl_chart_type(kind: Any, has_markers: bool) -> Optional[XL_CHART_TYPE]:
    if kind in _FIXED_XL_CHART_TYPE_BY_KIND:
        return _FIXED_XL_CHART_TYPE_BY_KIND[kind]
    if kind == "line":
        return XL_CHART_TYPE.LINE_MARKERS if has_markers else XL_CHART_TYPE.LINE
    if kind == "radar":
        return XL_CHART_TYPE.RADAR_MARKERS if has_markers else XL_CHART_TYPE.RADAR
    return None


def _hex_from_css_color(value: Any) -> Optional[str]:
    if not isinstance(value, str):
        return None
    text = value.strip()
    if text.startswith("#"):
        hex_part = text[1:]
        if len(hex_part) == 3:
            hex_part = "".join(ch * 2 for ch in hex_part)
        if len(hex_part) == 6 and all(
            ch in "0123456789abcdefABCDEF" for ch in hex_part
        ):
            return hex_part.upper()
        return None
    if text.startswith("rgb"):
        try:
            inner = text[text.index("(") + 1 : text.index(")")]
            parts = [p.strip() for p in inner.split(",")]
            r, g, b = (int(float(p)) for p in parts[:3])
        except (ValueError, IndexError):
            return None
        if all(0 <= c <= 255 for c in (r, g, b)):
            return f"{r:02X}{g:02X}{b:02X}"
        return None
    return None


def _alpha_from_css_color(value: Any) -> Optional[float]:
    """The 4th component of an `rgba(r, g, b, a)` string, or None for
    anything else (opaque `rgb()`/`#hex`, malformed input, or a=1). Kept
    deliberately separate from _hex_from_css_color rather than changing that
    function's return shape - several call sites (and its own pinned tests)
    depend on its existing "6-hex-or-None" contract. 8-digit #RRGGBBAA is not
    handled here because _hex_from_css_color already rejects any hex string
    that isn't 3 or 6 characters, so such a color never reaches a fill site
    to begin with; the two would need to change together if that ever does."""
    if not isinstance(value, str):
        return None
    text = value.strip()
    if not text.startswith("rgb"):
        return None
    try:
        inner = text[text.index("(") + 1 : text.index(")")]
        parts = [p.strip() for p in inner.split(",")]
        alpha = float(parts[3])
    except (ValueError, IndexError):
        return None
    if not (0.0 <= alpha < 1.0):
        # >=1.0 is opaque - nothing to inject; out-of-range is malformed, and
        # failing to today's opaque behavior beats emitting an invalid
        # ST_PositivePercentage into the file.
        return None
    return alpha


def _pt_from_css_px(value: Any) -> Optional[Length]:
    """Convert a captured CSS-px font size to a python-pptx `Length` in
    points, or None if the value is absent/invalid. Never raises - a garbage
    captured size (out of range, non-numeric, NaN/inf) must fall back to
    PowerPoint's own default rather than corrupt or crash the export."""
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        return None
    if value != value or value in (float("inf"), float("-inf")):  # NaN/inf
        return None
    if not (_MIN_FONT_PX <= value <= _MAX_FONT_PX):
        return None
    return Pt(value * _PT_PER_CSS_PX)


def _rect_to_emu(rect: dict, emu_per_px_x: float, emu_per_px_y: float) -> Rect:
    left = int(round(float(rect.get("left", 0)) * emu_per_px_x))
    top = int(round(float(rect.get("top", 0)) * emu_per_px_y))
    width = int(round(float(rect.get("width", 0)) * emu_per_px_x))
    height = int(round(float(rect.get("height", 0)) * emu_per_px_y))
    return left, top, width, height


def _shape_rect(shape) -> Rect:
    return shape.left, shape.top, shape.width, shape.height


def _iou(a: Rect, b: Rect) -> float:
    ax0, ay0, aw, ah = a
    bx0, by0, bw, bh = b
    if aw <= 0 or ah <= 0 or bw <= 0 or bh <= 0:
        return 0.0
    ax1, ay1 = ax0 + aw, ay0 + ah
    bx1, by1 = bx0 + bw, by0 + bh
    inter_w = max(0, min(ax1, bx1) - max(ax0, bx0))
    inter_h = max(0, min(ay1, by1) - max(ay0, by0))
    inter = inter_w * inter_h
    if inter <= 0:
        return 0.0
    union = aw * ah + bw * bh - inter
    if union <= 0:
        return 0.0
    return inter / union


def _best_overlap_match(
    target_rect: Rect, picture_shapes: list
) -> tuple[Optional[Any], float, float]:
    scored = sorted(
        ((_iou(target_rect, _shape_rect(shape)), shape) for shape in picture_shapes),
        key=lambda pair: pair[0],
        reverse=True,
    )
    if not scored:
        return None, 0.0, 0.0
    best_score, best_shape = scored[0]
    second_score = scored[1][0] if len(scored) > 1 else 0.0
    return best_shape, best_score, second_score


def _build_chart_data(captured_chart: dict) -> Optional[CategoryChartData]:
    labels = captured_chart.get("labels")
    datasets = captured_chart.get("datasets")
    if not isinstance(labels, list) or not labels:
        return None
    if not isinstance(datasets, list) or not datasets:
        return None

    data = CategoryChartData()
    data.categories = [str(label) for label in labels]
    for dataset in datasets:
        if not isinstance(dataset, dict):
            return None
        values = dataset.get("data")
        if not isinstance(values, list) or len(values) != len(labels):
            return None
        try:
            numeric_values = [float(v) for v in values]
        except (TypeError, ValueError):
            return None
        if not all(v == v and abs(v) != float("inf") for v in numeric_values):  # NaN/Inf guard
            return None
        data.add_series(str(dataset.get("label") or "Series"), numeric_values)
    return data


def _series_color(dataset: dict, index: int) -> Optional[str]:
    for field in ("backgroundColor", "borderColor"):
        color = dataset.get(field)
        if isinstance(color, list):
            candidate = color[index] if index < len(color) else (
                color[0] if color else None
            )
        else:
            candidate = color
        hex_color = _hex_from_css_color(candidate)
        if hex_color:
            return hex_color
    return None


def _series_raw_color(dataset: dict, index: int) -> Any:
    """Same field/index resolution as _series_color, but returns the raw
    candidate (whatever produced a truthy _hex_from_css_color) instead of the
    parsed hex - needed to also recover its alpha channel via
    _alpha_from_css_color, which _series_color's return value has already
    discarded. Deliberately a separate function rather than changing
    _series_color's return shape, which callers and tests already depend on."""
    for field in ("backgroundColor", "borderColor"):
        color = dataset.get(field)
        if isinstance(color, list):
            candidate = color[index] if index < len(color) else (
                color[0] if color else None
            )
        else:
            candidate = color
        if _hex_from_css_color(candidate):
            return candidate
    return None


def _no_fill_spPr() -> etree._Element:
    """A `<c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>` element."""
    spPr = etree.Element(qn("c:spPr"))
    etree.SubElement(spPr, qn("a:noFill"))
    ln = etree.SubElement(spPr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))
    return spPr


def _apply_fill_alpha(fill, alpha: float) -> None:
    """Make a solid `fill` translucent by appending `<a:alpha>` to its
    `<a:srgbClr>` element. python-pptx exposes no high-level API for fill
    transparency - `ColorFormat.rgb`'s setter only ever writes the color
    value, never an alpha transform. `<a:alpha>` is a member of
    EG_ColorTransform, an unbounded-repeat *choice* schema group, so
    appending it as the last child of `<a:srgbClr>` is schema-valid
    regardless of what other transforms (none, in this codebase) are already
    present - unlike _make_chart_background_transparent above, there is no
    fixed insertion position to get right here.

    MUST be called after `fill.solid(); fill.fore_color.rgb = ...`, never
    before: ColorFormat.rgb's setter (pptx/dml/color.py) calls
    get_or_change_to_srgbClr(), which replaces the color element outright: an
    alpha written first would be discarded when the color is then assigned.
    """
    xFill = fill.fore_color._xFill
    srgbClr = xFill.find(qn("a:srgbClr"))
    if srgbClr is None:
        return
    for existing in srgbClr.findall(qn("a:alpha")):
        srgbClr.remove(existing)
    alpha_el = etree.SubElement(srgbClr, qn("a:alpha"))
    alpha_el.set("val", str(int(round(alpha * 100000))))


def _make_chart_background_transparent(chart) -> None:
    """
    Native PPTX charts default to an opaque white chart-area and plot-area
    fill (PowerPoint's theme default), unlike the original Chart.js canvas,
    which is transparent and lets the surrounding card's own background
    (often dark) show through - the "white box" mismatch reported after
    export. python-pptx exposes no high-level API for chart/plot-area fill,
    so this inserts `<c:spPr>` directly at the two schema-mandated
    positions: as the last child of `c:plotArea` (per CT_PlotArea), and
    immediately after `c:chart` inside `c:chartSpace` (per CT_ChartSpace's
    `_tag_seq`, before `c:txPr`/`c:externalData` if present) - both
    confirmed against python-pptx's own oxml schema classes and verified by
    round-tripping a real file through save/reload.
    """
    chartSpace = chart._chartSpace
    chart_el = chartSpace.find(qn("c:chart"))
    plotArea = chart_el.find(qn("c:plotArea"))

    if plotArea.find(qn("c:spPr")) is None:
        plotArea.append(_no_fill_spPr())
    if chart_el.find(qn("c:spPr")) is None:
        chart_el.addnext(_no_fill_spPr())


def _apply_chart_styling(
    chart, xl_type: XL_CHART_TYPE, kind: str, captured_chart: dict
) -> None:
    # Resolved once, up front, so every site below can share them.
    # default_pt backs chart.font.size, the chart-wide baseline every text
    # element inherits unless it's overridden below - this single value is
    # what actually fixes native text defaulting to PowerPoint's 18pt instead
    # of the deck's real ~10-11pt, for every element this function doesn't
    # otherwise touch.
    data_label_pt = _pt_from_css_px(captured_chart.get("dataLabelFontSize"))
    tick_pt = _pt_from_css_px(captured_chart.get("tickFontSize"))
    legend_pt = _pt_from_css_px(captured_chart.get("legendFontSize"))
    title_pt = _pt_from_css_px(captured_chart.get("titleFontSize"))
    axis_title_pt = _pt_from_css_px(captured_chart.get("axisTitleFontSize"))
    default_pt = next(
        (pt for pt in (tick_pt, data_label_pt, legend_pt) if pt is not None), None
    )
    if default_pt is not None:
        chart.font.size = default_pt

    legend_cfg = captured_chart.get("legend") or {}
    chart.has_legend = bool(legend_cfg.get("show"))
    if chart.has_legend:
        legend_position = legend_cfg.get("position")
        chart.legend.position = _XL_LEGEND_POSITION_BY_NAME.get(
            legend_position, XL_LEGEND_POSITION.BOTTOM
        )
        chart.legend.include_in_layout = False
        legend_color = _hex_from_css_color(captured_chart.get("legendColor"))
        if legend_color:
            chart.legend.font.color.rgb = RGBColor.from_string(legend_color)
        legend_size = next(
            (pt for pt in (legend_pt, default_pt) if pt is not None), None
        )
        if legend_size is not None:
            chart.legend.font.size = legend_size

    title = captured_chart.get("title")
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = str(title)
        title_font = chart.chart_title.text_frame.paragraphs[0].font
        title_color = _hex_from_css_color(captured_chart.get("titleColor"))
        if title_color:
            title_font.color.rgb = RGBColor.from_string(title_color)
        if title_pt is not None:
            title_font.size = title_pt

    is_line_family = xl_type in (
        XL_CHART_TYPE.LINE,
        XL_CHART_TYPE.LINE_MARKERS,
        XL_CHART_TYPE.RADAR,
        XL_CHART_TYPE.RADAR_MARKERS,
    )
    datasets = captured_chart.get("datasets") or []
    plot = chart.plots[0]

    for series_index, series in enumerate(plot.series):
        if series_index >= len(datasets) or not isinstance(datasets[series_index], dict):
            continue
        dataset = datasets[series_index]

        point_colors = None
        if kind in _POINT_COLOR_KINDS:
            for field in ("backgroundColor", "borderColor"):
                value = dataset.get(field)
                if isinstance(value, list) and value:
                    point_colors = value
                    break

        if point_colors is not None:
            for point_index, point in enumerate(series.points):
                raw = (
                    point_colors[point_index]
                    if point_index < len(point_colors)
                    else None
                )
                color = _hex_from_css_color(raw)
                if not color:
                    continue
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor.from_string(color)
                alpha = _alpha_from_css_color(raw)
                if alpha is not None:
                    _apply_fill_alpha(point.format.fill, alpha)
            continue

        color = _series_color(dataset, series_index)
        if not color:
            continue
        if is_line_family:
            series.format.line.color.rgb = RGBColor.from_string(color)
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(color)
            alpha = _alpha_from_css_color(_series_raw_color(dataset, series_index))
            if alpha is not None:
                _apply_fill_alpha(series.format.fill, alpha)

    if captured_chart.get("hasDataLabels"):
        plot.has_data_labels = True
        data_labels = plot.data_labels
        if kind in ("pie", "donut"):
            # show_percentage (PowerPoint's own computed share-of-whole) is
            # the semantically correct label for any pie/donut regardless of
            # whether the source values already happen to be percentages.
            data_labels.show_percentage = True
            data_labels.show_value = False
            if kind in _OUTSIDE_END_LABEL_KINDS:
                # A bare percentage on a pie slice doesn't say what it is
                # without cross-referencing the legend - show the category
                # name alongside it. OUTSIDE_END positioning (needed to fit
                # "Category NN%" outside a slim slice) is only legal for
                # pie, not doughnut - see _OUTSIDE_END_LABEL_KINDS. Donut
                # labels sit inside the ring at PowerPoint's own default
                # position instead, so show_category_name stays off there to
                # avoid crowding.
                data_labels.show_category_name = True
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            else:
                data_labels.show_category_name = False
        else:
            # Nothing else ever sets a show_* flag for non-circular charts,
            # so plot.has_data_labels = True alone emits a <c:dLbls> whose
            # <c:showVal> defaults to "0" - a data-label block that displays
            # literally nothing. Confirmed in a real export: every bar/line/
            # area/radar chart's data labels were invisible for this reason.
            data_labels.show_value = True
        data_label_color = _hex_from_css_color(
            captured_chart.get("dataLabelColor") or captured_chart.get("tickColor")
        )
        if data_label_color:
            data_labels.font.color.rgb = RGBColor.from_string(data_label_color)
        data_label_size = next(
            (pt for pt in (data_label_pt, default_pt) if pt is not None), None
        )
        if data_label_size is not None:
            data_labels.font.size = data_label_size

    font_family = captured_chart.get("fontFamily")
    if isinstance(font_family, str) and font_family.strip():
        chart.font.name = font_family.strip()

    if kind in _REVERSE_CATEGORY_ORDER_KINDS:
        chart.category_axis.reverse_order = True

    if kind not in _NO_AXIS_KINDS:
        # tick_color and tick_pt are independent siblings on purpose (an
        # earlier version nested the size write inside `if tick_color:`,
        # which meant a chart with a captured size but no captured color -
        # or vice versa - silently lost whichever wasn't present).
        cat_tick_labels = chart.category_axis.tick_labels
        val_tick_labels = chart.value_axis.tick_labels
        tick_color = _hex_from_css_color(captured_chart.get("tickColor"))
        if tick_color:
            cat_tick_labels.font.color.rgb = RGBColor.from_string(tick_color)
            val_tick_labels.font.color.rgb = RGBColor.from_string(tick_color)
        if tick_pt is not None:
            cat_tick_labels.font.size = tick_pt
            val_tick_labels.font.size = tick_pt

    if kind in _AXIS_TITLE_SUPPORTED_KINDS:
        axis_titles = captured_chart.get("axisTitles") or {}
        x_title, y_title = axis_titles.get("x"), axis_titles.get("y")
        axis_title_color = _hex_from_css_color(captured_chart.get("axisTitleColor"))
        if x_title:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = str(x_title)
            x_title_font = chart.category_axis.axis_title.text_frame.paragraphs[0].font
            if axis_title_color:
                x_title_font.color.rgb = RGBColor.from_string(axis_title_color)
            if axis_title_pt is not None:
                x_title_font.size = axis_title_pt
        if y_title:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = str(y_title)
            y_title_font = chart.value_axis.axis_title.text_frame.paragraphs[0].font
            if axis_title_color:
                y_title_font.color.rgb = RGBColor.from_string(axis_title_color)
            if axis_title_pt is not None:
                y_title_font.size = axis_title_pt


def _replace_picture_with_native_chart(
    slide, picture_shape, xl_type: XL_CHART_TYPE, kind: str, captured_chart: dict
) -> bool:
    try:
        chart_data = _build_chart_data(captured_chart)
        if chart_data is None:
            return False

        sp_tree = slide.shapes._spTree
        picture_element = picture_shape._element
        original_index = list(sp_tree).index(picture_element)
        left, top, width, height = _shape_rect(picture_shape)

        # The picture must come out of the tree before add_chart() can be
        # called (python-pptx has no "build a chart, then swap it in" API),
        # but add_chart() itself can raise - a malformed category/series
        # shape, for instance. An earlier version removed the picture here
        # and only caught the failure in the outer except below, which
        # returned False (chart upgrade skipped) while leaving the picture
        # already gone from the tree - so a single failing chart in an
        # otherwise-successful deck saved with a blank hole where a chart
        # used to be, and the log line ("failed to build native chart for
        # one slide") actively suggested the flattened image was still
        # there. Restoring the picture on failure keeps the documented
        # contract ("upgrade fails -> chart stays a flattened image") true
        # in the one case that used to violate it.
        sp_tree.remove(picture_element)
        try:
            graphic_frame = slide.shapes.add_chart(
                xl_type, left, top, width, height, chart_data
            )
        except Exception:
            sp_tree.insert(original_index, picture_element)
            raise

        # add_chart() appends at the end of the tree; move it back to the
        # picture's original position so z-order/stacking is unchanged.
        sp_tree.remove(graphic_frame._element)
        sp_tree.insert(original_index, graphic_frame._element)

        try:
            _make_chart_background_transparent(graphic_frame.chart)
            _apply_chart_styling(graphic_frame.chart, xl_type, kind, captured_chart)
        except Exception:
            LOGGER.exception(
                "pptx_native_chart_service: styling failed, chart kept with defaults"
            )
        return True
    except Exception:
        LOGGER.exception(
            "pptx_native_chart_service: failed to build native chart for one slide"
        )
        return False


def _try_upgrade_one_chart(
    slide,
    captured_chart: dict,
    emu_per_px_x: float,
    emu_per_px_y: float,
    claimed_shape_ids: set,
) -> bool:
    kind = captured_chart.get("kind")
    has_markers = bool(captured_chart.get("hasMarkers"))
    xl_type = _resolve_xl_chart_type(kind, has_markers)
    if xl_type is None:
        LOGGER.info(
            "pptx_native_chart_service: skip kind=%s (unsupported, staying a picture)",
            kind,
        )
        return False

    bounding_box = captured_chart.get("boundingBox")
    if not isinstance(bounding_box, dict):
        LOGGER.info(
            "pptx_native_chart_service: skip kind=%s (missing boundingBox)", kind
        )
        return False
    target_rect = _rect_to_emu(bounding_box, emu_per_px_x, emu_per_px_y)
    if target_rect[2] <= 0 or target_rect[3] <= 0:
        LOGGER.info(
            "pptx_native_chart_service: skip kind=%s (degenerate boundingBox %s)",
            kind,
            bounding_box,
        )
        return False

    picture_shapes = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and id(shape._element) not in claimed_shape_ids
    ]
    best_shape, best_score, second_score = _best_overlap_match(
        target_rect, picture_shapes
    )
    if best_shape is None:
        LOGGER.info(
            "pptx_native_chart_service: skip kind=%s (no picture shapes on slide)",
            kind,
        )
        return False
    if best_score < _IOU_MIN or (best_score - second_score) < _IOU_AMBIGUITY_MARGIN:
        LOGGER.info(
            "pptx_native_chart_service: skip kind=%s (no confident geometry match: "
            "best_iou=%.3f second_best_iou=%.3f candidates=%s)",
            kind,
            best_score,
            second_score,
            len(picture_shapes),
        )
        return False

    claimed_shape_ids.add(id(best_shape._element))
    upgraded = _replace_picture_with_native_chart(
        slide, best_shape, xl_type, kind, captured_chart
    )
    LOGGER.info(
        "pptx_native_chart_service: %s kind=%s xl_type=%s iou=%.3f",
        "upgraded" if upgraded else "failed to build native chart for",
        kind,
        xl_type,
        best_score,
    )
    return upgraded


def _group_by_slide_order(charts: list) -> dict[int, list[dict]]:
    grouped: dict[int, list[dict]] = {}
    for chart in charts:
        if not isinstance(chart, dict):
            continue
        index = chart.get("slideOrderIndex")
        if not isinstance(index, int) or isinstance(index, bool):
            continue
        grouped.setdefault(index, []).append(chart)
    return grouped


async def _expected_slide_count(presentation_id: uuid.UUID) -> int:
    async with async_session_maker() as session:
        result = await session.execute(
            select(SlideModel.id).where(SlideModel.presentation == presentation_id)
        )
        return len(result.all())


async def upgrade_flattened_charts_to_native(
    pptx_path: str, token: Optional[str], presentation_id: uuid.UUID
) -> None:
    """Best-effort. Never raises. Leaves pptx_path untouched on any failure,
    ambiguity, or absence of captured chart data."""
    if not token:
        LOGGER.info(
            "pptx_native_chart_service: no capture token for this export "
            "(pdf export, or pptx export minted none) - skipping"
        )
        return
    try:
        capture = chart_capture_store.take_capture(token)
        if not capture:
            LOGGER.info(
                "pptx_native_chart_service: no capture found for token=%s "
                "(the export page's report never arrived, or arrived after "
                "this export finished) - keeping flattened images",
                token,
            )
            return
        charts = capture.get("charts")
        if not isinstance(charts, list) or not charts:
            LOGGER.info(
                "pptx_native_chart_service: capture for token=%s had zero charts "
                "- keeping flattened images",
                token,
            )
            return

        charts_by_slide = _group_by_slide_order(charts)
        if not charts_by_slide:
            LOGGER.info(
                "pptx_native_chart_service: capture for token=%s had %s chart(s) "
                "but none had a usable slideOrderIndex - keeping flattened images",
                token,
                len(charts),
            )
            return

        expected_count = await _expected_slide_count(presentation_id)

        prs = Presentation(pptx_path)
        if expected_count and len(prs.slides) != expected_count:
            LOGGER.warning(
                "pptx_native_chart_service: slide count mismatch (pptx=%s, expected=%s), "
                "skipping native chart upgrade",
                len(prs.slides),
                expected_count,
            )
            return

        LOGGER.info(
            "pptx_native_chart_service: token=%s captured %s chart(s) across %s slide(s), "
            "pptx has %s slide(s) - starting native chart matching",
            token,
            len(charts),
            len(charts_by_slide),
            len(prs.slides),
        )

        emu_per_px_x = prs.slide_width / _DESIGN_WIDTH_PX
        emu_per_px_y = prs.slide_height / _DESIGN_HEIGHT_PX

        upgraded_count = 0
        for slide_order_index, slide in enumerate(prs.slides):
            slide_charts = charts_by_slide.get(slide_order_index)
            if not slide_charts:
                continue
            claimed_shape_ids: set = set()
            for captured_chart in slide_charts:
                if _try_upgrade_one_chart(
                    slide,
                    captured_chart,
                    emu_per_px_x,
                    emu_per_px_y,
                    claimed_shape_ids,
                ):
                    upgraded_count += 1

        if upgraded_count:
            prs.save(pptx_path)
        LOGGER.info(
            "pptx_native_chart_service: upgraded %s of %s captured chart(s) to native PPTX charts",
            upgraded_count,
            len(charts),
        )
    except Exception:
        LOGGER.exception(
            "pptx_native_chart_service: native chart upgrade failed, keeping flattened images"
        )
