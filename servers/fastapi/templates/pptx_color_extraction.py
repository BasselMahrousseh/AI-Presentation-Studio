"""Extract a representative brand color palette from an uploaded PPTX file.

Used to let a Smart-mode e& generation follow a user-supplied brand deck's
real colors instead of the hardcoded default palette (see
utils/smart_brand_templates.py). Two signals are combined because either one
alone is unreliable in practice:

- Theme colors (`<a:clrScheme>` in the deck's theme XML) are the "correct"
  PowerPoint concept of a brand palette, but a great many real decks
  (including every deck this app itself exports) never touch them and just
  leave the default Office theme colors in place, so relying on theme colors
  alone silently returns a generic blue/orange palette that has nothing to do
  with the deck's actual look.
- Scanning every shape's literal fill/line/font color and ranking by
  frequency reliably surfaces the colors actually used on the page,
  regardless of whether the author used theme slots or direct colors.
"""

from __future__ import annotations

import asyncio
import tempfile
from collections import Counter
from pathlib import Path
from typing import Iterable

from fastapi import File, HTTPException, UploadFile
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel

from templates.fonts_and_slides_preview import (
    _raise_if_font_check_upload_too_large,
    _validate_pptx_package,
    _write_bytes_to_path,
)

_THEME_RELATIONSHIP_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
)
_DRAWINGML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
_THEME_ACCENT_SLOTS = ("accent1", "accent2", "accent3", "accent4", "accent5", "accent6")
_THEME_DARK_SLOTS = ("dk2", "dk1")
_NEAR_NEUTRAL_CHANNEL_SPREAD = 18
_MAX_BRAND_COLORS = 4
_MAX_SLIDES_SCANNED = 60


def _is_near_neutral(hex_color: str) -> bool:
    """True for near-white/near-black/low-saturation grays, which are almost
    always canvas/text colors rather than brand accents."""
    try:
        r, g, b = (int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
    except (ValueError, IndexError):
        return True
    return max(r, g, b) - min(r, g, b) < _NEAR_NEUTRAL_CHANNEL_SPREAD


def _theme_accent_colors(prs: Presentation) -> list[str]:
    colors: list[str] = []
    seen_partnames: set[str] = set()
    for master in prs.slide_masters:
        try:
            theme_part = master.part.part_related_by(_THEME_RELATIONSHIP_TYPE)
        except KeyError:
            continue
        partname = str(theme_part.partname)
        if partname in seen_partnames:
            continue
        seen_partnames.add(partname)
        try:
            theme_xml = etree.fromstring(theme_part.blob)
        except etree.XMLSyntaxError:
            continue
        scheme = theme_xml.find(".//a:clrScheme", _DRAWINGML_NS)
        if scheme is None:
            continue
        slots = {child.tag.rsplit("}", 1)[-1]: child for child in scheme}
        for slot_name in (*_THEME_DARK_SLOTS, *_THEME_ACCENT_SLOTS):
            slot = slots.get(slot_name)
            if slot is None:
                continue
            srgb = slot.find("a:srgbClr", _DRAWINGML_NS)
            sys_clr = slot.find("a:sysClr", _DRAWINGML_NS)
            value = srgb.get("val") if srgb is not None else (
                sys_clr.get("lastClr") if sys_clr is not None else None
            )
            if value:
                colors.append(value.upper())
    return colors


def _iter_shape_colors(shape) -> Iterable[str]:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                yield from _iter_shape_colors(sub_shape)
            return
    except (AttributeError, ValueError):
        pass

    for attribute_name in ("fill", "line"):
        try:
            fill_format = getattr(shape, attribute_name, None)
            if fill_format is None:
                continue
            color = fill_format.fore_color
            if color.type is not None and color.rgb is not None:
                yield str(color.rgb).upper()
        except (AttributeError, KeyError, ValueError, TypeError):
            continue

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    color = run.font.color
                    if color.type is not None and color.rgb is not None:
                        yield str(color.rgb).upper()
                except (AttributeError, KeyError, ValueError, TypeError):
                    continue


def _shape_fill_color_counts(prs: Presentation) -> Counter:
    counts: Counter = Counter()
    for slide in list(prs.slides)[:_MAX_SLIDES_SCANNED]:
        for shape in slide.shapes:
            for hex_color in _iter_shape_colors(shape):
                counts[hex_color] += 1
    return counts


def extract_pptx_brand_colors(
    pptx_path: str, max_colors: int = _MAX_BRAND_COLORS
) -> list[str]:
    """Return up to `max_colors` distinct, non-neutral hex colors (no `#`
    prefix, uppercase), ranked by how prominently they're actually used in
    the deck. Falls back to theme accent colors only to break ties / fill
    remaining slots, since they're a weaker signal (see module docstring)."""
    prs = Presentation(pptx_path)
    shape_counts = _shape_fill_color_counts(prs)
    ranked_by_usage = [
        color
        for color, _count in shape_counts.most_common()
        if not _is_near_neutral(color)
    ]

    theme_colors = [
        color for color in _theme_accent_colors(prs) if not _is_near_neutral(color)
    ]

    ordered_candidates: list[str] = []
    for color in (*ranked_by_usage, *theme_colors):
        if color not in ordered_candidates:
            ordered_candidates.append(color)

    return ordered_candidates[:max_colors]


class PptxColorPaletteResponse(BaseModel):
    colors: list[str]


INVALID_COLOR_PALETTE_PPTX_ERROR = "Invalid file type. Expected PPTX file"


async def extract_pptx_color_palette_handler(
    pptx_file: UploadFile = File(..., description="PPTX file to extract a color palette from"),
) -> PptxColorPaletteResponse:
    filename = getattr(pptx_file, "filename", "") or ""
    if not filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail=INVALID_COLOR_PALETTE_PPTX_ERROR)
    _raise_if_font_check_upload_too_large(getattr(pptx_file, "size", None))

    with tempfile.TemporaryDirectory() as temp_dir:
        pptx_path = str(Path(temp_dir) / "presentation.pptx")
        pptx_content = await pptx_file.read()
        _raise_if_font_check_upload_too_large(len(pptx_content))
        await asyncio.to_thread(_write_bytes_to_path, pptx_path, pptx_content)
        await asyncio.to_thread(_validate_pptx_package, pptx_path)
        hex_colors = await asyncio.to_thread(extract_pptx_brand_colors, pptx_path)

    return PptxColorPaletteResponse(colors=[f"#{color}" for color in hex_colors])
