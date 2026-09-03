import asyncio
import io
import uuid

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from services import pptx_native_chart_service as svc


# ---------------------------------------------------------------------------
# Chart-type mapping
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "kind,has_markers,expected",
    [
        ("bar", False, XL_CHART_TYPE.COLUMN_CLUSTERED),
        ("horizontal_bar", False, XL_CHART_TYPE.BAR_CLUSTERED),
        ("stacked_bar", False, XL_CHART_TYPE.COLUMN_STACKED),
        ("horizontal_stacked_bar", False, XL_CHART_TYPE.BAR_STACKED),
        ("pie", False, XL_CHART_TYPE.PIE),
        ("donut", False, XL_CHART_TYPE.DOUGHNUT),
        ("area", False, XL_CHART_TYPE.AREA),
        ("line", False, XL_CHART_TYPE.LINE),
        ("line", True, XL_CHART_TYPE.LINE_MARKERS),
        ("radar", False, XL_CHART_TYPE.RADAR),
        ("radar", True, XL_CHART_TYPE.RADAR_MARKERS),
    ],
)
def test_resolve_xl_chart_type_supported_kinds(kind, has_markers, expected):
    assert svc._resolve_xl_chart_type(kind, has_markers) == expected


@pytest.mark.parametrize("kind", ["polar_area", "scatter", "bubble", "unknown", None, 42])
def test_resolve_xl_chart_type_unsupported_kinds_return_none(kind):
    assert svc._resolve_xl_chart_type(kind, has_markers=False) is None


# ---------------------------------------------------------------------------
# Color parsing
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "value,expected",
    [
        ("#ff0000", "FF0000"),
        ("#F00", "FF0000"),
        ("rgb(0, 128, 255)", "0080FF"),
        ("rgba(10, 20, 30, 0.5)", "0A141E"),
        ("not-a-color", None),
        (None, None),
        (123, None),
        ("#12345", None),
    ],
)
def test_hex_from_css_color(value, expected):
    assert svc._hex_from_css_color(value) == expected


# ---------------------------------------------------------------------------
# Geometry matching
# ---------------------------------------------------------------------------


def test_iou_exact_overlap_is_one():
    rect = (0, 0, 100, 100)
    assert svc._iou(rect, rect) == pytest.approx(1.0)


def test_iou_no_overlap_is_zero():
    assert svc._iou((0, 0, 10, 10), (100, 100, 10, 10)) == 0.0


def test_iou_partial_overlap():
    a = (0, 0, 100, 100)
    b = (50, 0, 100, 100)
    # intersection 50x100=5000, union 100*100*2-5000=15000
    assert svc._iou(a, b) == pytest.approx(5000 / 15000)


class _FakeShape:
    def __init__(self, left, top, width, height):
        self.left, self.top, self.width, self.height = left, top, width, height
        self._element = object()


def test_best_overlap_match_picks_highest_scoring_confident_match():
    target = (0, 0, 100, 100)
    exact = _FakeShape(0, 0, 100, 100)
    far = _FakeShape(1000, 1000, 100, 100)
    best, best_score, second_score = svc._best_overlap_match(target, [far, exact])
    assert best is exact
    assert best_score == pytest.approx(1.0)
    assert second_score == 0.0


def test_best_overlap_match_no_candidates_returns_none():
    best, best_score, second_score = svc._best_overlap_match((0, 0, 10, 10), [])
    assert best is None
    assert best_score == 0.0
    assert second_score == 0.0


# ---------------------------------------------------------------------------
# Chart data assembly
# ---------------------------------------------------------------------------


def test_build_chart_data_happy_path():
    chart_data = svc._build_chart_data(
        {
            "labels": ["a", "b", "c"],
            "datasets": [{"label": "S1", "data": [1, 2, 3]}],
        }
    )
    assert chart_data is not None
    assert [c.label for c in chart_data.categories] == ["a", "b", "c"]


@pytest.mark.parametrize(
    "payload",
    [
        {"labels": [], "datasets": [{"label": "S1", "data": []}]},
        {"labels": ["a", "b"], "datasets": []},
        {"labels": ["a", "b"], "datasets": [{"label": "S1", "data": [1]}]},  # length mismatch
        {"labels": ["a", "b"], "datasets": [{"label": "S1", "data": ["x", "y"]}]},  # non-numeric
        {"labels": ["a", "b"], "datasets": [{"label": "S1", "data": [1, float("nan")]}]},
    ],
)
def test_build_chart_data_rejects_malformed_input(payload):
    assert svc._build_chart_data(payload) is None


# ---------------------------------------------------------------------------
# Grouping by slide order
# ---------------------------------------------------------------------------


def test_group_by_slide_order():
    charts = [
        {"slideOrderIndex": 0, "kind": "bar"},
        {"slideOrderIndex": 0, "kind": "line"},
        {"slideOrderIndex": 2, "kind": "pie"},
        {"slideOrderIndex": "not-an-int", "kind": "bad"},
        "not-a-dict",
    ]
    grouped = svc._group_by_slide_order(charts)
    assert set(grouped.keys()) == {0, 2}
    assert len(grouped[0]) == 2
    assert len(grouped[2]) == 1


# ---------------------------------------------------------------------------
# End-to-end python-pptx object assembly
# ---------------------------------------------------------------------------


def _make_presentation_with_picture(left_px, top_px, width_px, height_px):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    emu_x = prs.slide_width / 1280
    emu_y = prs.slide_height / 720

    img = Image.new("RGB", (10, 10), color="red")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    picture = slide.shapes.add_picture(
        buf,
        Emu(int(left_px * emu_x)),
        Emu(int(top_px * emu_y)),
        Emu(int(width_px * emu_x)),
        Emu(int(height_px * emu_y)),
    )
    return prs, slide, picture, emu_x, emu_y


def test_try_upgrade_one_chart_replaces_matching_picture_with_native_chart():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(100, 100, 400, 300)
    captured_chart = {
        "kind": "bar",
        "hasMarkers": False,
        "boundingBox": {"left": 100, "top": 100, "width": 400, "height": 300},
        "labels": ["Q1", "Q2", "Q3"],
        "datasets": [{"label": "Revenue", "data": [10, 20, 30], "backgroundColor": "#3366FF"}],
        "title": "Quarterly Revenue",
        "legend": {"show": True},
        "axisTitles": {"x": "Quarter", "y": "USD"},
    }
    claimed: set = set()
    upgraded = svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed)
    assert upgraded is True

    shapes = list(slide.shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.CHART
    assert shapes[0].chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert list(shapes[0].chart.plots[0].categories) == ["Q1", "Q2", "Q3"]


def test_failed_add_chart_leaves_the_original_picture_in_place(monkeypatch):
    """A slide's picture must not disappear when the native-chart upgrade
    fails partway through. Reproduces a real bug: add_chart() was called
    after the picture had already been removed from the shape tree, and the
    failure path returned False without restoring it - so a deck with one
    bad chart among several good ones still saved, with a blank hole where
    that chart used to be, while the log line said only "failed to build
    native chart for one slide" (which reads as "the picture is still
    there"). This is the one chart in the deck whose real slide.shapes call
    must raise, so the fixture patches add_chart on that specific
    SlideShapes instance rather than monkeypatching python-pptx globally."""
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(
        100, 100, 400, 300
    )
    original_element = picture._element
    original_index = list(slide.shapes._spTree).index(original_element)

    def _boom(*args, **kwargs):
        raise RuntimeError("malformed chart data")

    monkeypatch.setattr(slide.shapes, "add_chart", _boom)

    captured_chart = {
        "kind": "bar",
        "hasMarkers": False,
        "boundingBox": {"left": 100, "top": 100, "width": 400, "height": 300},
        "labels": ["Q1", "Q2", "Q3"],
        "datasets": [{"label": "Revenue", "data": [10, 20, 30]}],
    }
    claimed: set = set()
    upgraded = svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed)

    assert upgraded is False
    shapes = list(slide.shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
    # The exact original element, not merely "a picture" - proves it was
    # restored rather than a new one being created.
    assert shapes[0]._element is original_element
    # Position/order preserved too: it must land back at its original spot
    # in the shape tree, not appended at the end.
    sp_tree = slide.shapes._spTree
    assert list(sp_tree).index(original_element) == original_index


# ---------------------------------------------------------------------------
# Visual polish: per-point colors, data labels, font, category axis order
# ---------------------------------------------------------------------------


def test_per_point_colors_applied_to_single_series_bar_chart():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["Tech", "Retail", "Education"],
        "datasets": [
            {
                "label": "Adoption",
                "data": [78, 44, 33],
                "backgroundColor": ["#1ABC9C", "#7B68EE", "#7B68EE"],
            }
        ],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    chart = list(slide.shapes)[0].chart
    points = list(chart.plots[0].series[0].points)
    assert len(points) == 3
    assert points[0].format.fill.fore_color.rgb == RGBColor.from_string("1ABC9C")
    assert points[1].format.fill.fore_color.rgb == RGBColor.from_string("7B68EE")


def test_uniform_series_color_used_when_no_color_array_present():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2], "backgroundColor": "#3366FF"}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    chart = list(slide.shapes)[0].chart
    series = chart.plots[0].series[0]
    assert series.format.fill.fore_color.rgb == RGBColor.from_string("3366FF")


def test_line_chart_keeps_single_series_color_even_with_color_array():
    # Per-point coloring is only applied for the discrete/categorical kinds
    # (_POINT_COLOR_KINDS) - a line chart should keep one color for its
    # whole series even if backgroundColor happens to be an array.
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "line",
        "hasMarkers": True,
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [
            {
                "label": "S1",
                "data": [1, 2],
                "backgroundColor": ["#111111", "#222222"],
                "borderColor": "#3366FF",
            }
        ],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    chart = list(slide.shapes)[0].chart
    series = chart.plots[0].series[0]
    # Falls through to _series_color, which prefers backgroundColor[0] since
    # it's a list - confirms no per-point branch was taken (no exception from
    # treating a 2-point series as if it needed per-point coloring twice).
    assert series.format.line.color.rgb == RGBColor.from_string("111111")


def test_data_labels_enabled_when_captured():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "hasDataLabels": True,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.plots[0].has_data_labels is True


def test_data_labels_left_off_when_not_captured():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.plots[0].has_data_labels is False


def test_font_family_applied_when_captured():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "fontFamily": "Inter",
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.font.name == "Inter"


@pytest.mark.parametrize("kind", ["horizontal_bar", "horizontal_stacked_bar"])
def test_category_axis_reversed_for_horizontal_bar_kinds(kind):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": kind,
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.category_axis.reverse_order is True


def test_category_axis_not_reversed_for_vertical_bar():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.category_axis.reverse_order is False


def test_tick_and_data_label_colors_applied():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "hasDataLabels": True,
        "tickColor": "#AFC3CE",
        "dataLabelColor": "#FFFFFF",
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.category_axis.tick_labels.font.color.rgb == RGBColor.from_string("AFC3CE")
    assert chart.value_axis.tick_labels.font.color.rgb == RGBColor.from_string("AFC3CE")
    assert chart.plots[0].data_labels.font.color.rgb == RGBColor.from_string("FFFFFF")


def test_data_label_color_falls_back_to_tick_color_when_unset():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "hasDataLabels": True,
        "tickColor": "#AFC3CE",
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.plots[0].data_labels.font.color.rgb == RGBColor.from_string("AFC3CE")


def test_tick_color_skipped_for_pie_and_donut_without_error():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "donut",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2], "backgroundColor": ["#111111", "#222222"]}],
        "tickColor": "#AFC3CE",
    }
    claimed: set = set()
    # Must not raise even though donut has no category/value axis to color.
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True


def test_legend_and_title_colors_applied():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "legend": {"show": True},
        "legendColor": "#FFFFFF",
        "title": "My Chart",
        "titleColor": "#FFFFFF",
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.legend.font.color.rgb == RGBColor.from_string("FFFFFF")
    assert chart.chart_title.text_frame.paragraphs[0].font.color.rgb == RGBColor.from_string(
        "FFFFFF"
    )


# ---------------------------------------------------------------------------
# Font sizes (Fix 3)
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "px,expected_pt",
    [(13, 9.75), (14, 10.5), (18, 13.5), (1.0, 0.75), (200, 150.0)],
)
def test_pt_from_css_px(px, expected_pt):
    result = svc._pt_from_css_px(px)
    assert result is not None
    assert result.pt == pytest.approx(expected_pt)


@pytest.mark.parametrize(
    "value",
    [None, "14", True, False, float("nan"), float("inf"), float("-inf"), 0, -3, 5000, {}],
)
def test_pt_from_css_px_rejects_invalid_input(value):
    assert svc._pt_from_css_px(value) is None


def test_font_sizes_applied_from_captured_px():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "legend": {"show": True},
        "legendFontSize": 15,
        "title": "My Chart",
        "titleFontSize": 18,
        "hasDataLabels": True,
        "dataLabelFontSize": 14,
        "tickFontSize": 13,
        "axisTitles": {"x": "Category", "y": "Value"},
        "axisTitleFontSize": 13,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart

    assert chart.legend.font.size.pt == pytest.approx(11.25)
    assert chart.chart_title.text_frame.paragraphs[0].font.size.pt == pytest.approx(13.5)
    assert chart.plots[0].data_labels.font.size.pt == pytest.approx(10.5)
    assert chart.category_axis.tick_labels.font.size.pt == pytest.approx(9.75)
    assert chart.value_axis.tick_labels.font.size.pt == pytest.approx(9.75)
    assert chart.category_axis.axis_title.text_frame.paragraphs[
        0
    ].font.size.pt == pytest.approx(9.75)
    assert chart.value_axis.axis_title.text_frame.paragraphs[
        0
    ].font.size.pt == pytest.approx(9.75)


def test_font_sizes_absent_leaves_powerpoint_defaults():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "legend": {"show": True},
        "title": "My Chart",
        "hasDataLabels": True,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart

    assert chart.legend.font.size is None
    assert chart.chart_title.text_frame.paragraphs[0].font.size is None
    assert chart.plots[0].data_labels.font.size is None
    assert chart.category_axis.tick_labels.font.size is None
    # python-pptx's own chart-creation template always bakes exactly one
    # chart-wide `sz="1800"` into c:chartSpace/c:txPr (confirmed by building
    # a chart via add_chart() with zero styling calls and counting matches),
    # independent of anything this service does - only chart.font.size
    # legitimately reads that one back. Assert no *additional* sz attributes
    # were introduced anywhere else, rather than asserting none exist at all.
    assert chart._chartSpace.xml.count('sz="') == 1


def test_tick_font_size_applied_without_tick_color():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "tickFontSize": 13,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.category_axis.tick_labels.font.size.pt == pytest.approx(9.75)
    assert chart.category_axis.tick_labels.font.color.type is None


def test_tick_color_applied_without_tick_font_size():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "tickColor": "#AFC3CE",
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.category_axis.tick_labels.font.color.rgb == RGBColor.from_string("AFC3CE")
    assert chart.category_axis.tick_labels.font.size is None


def test_font_size_on_donut_does_not_raise():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = _donut_capture(tickFontSize=13)
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True


def test_chart_wide_default_font_size_from_tick_size():
    """When nothing overrides it, the chart-wide default (backing every text
    element this function doesn't otherwise touch) comes from tick size."""
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "tickFontSize": 13,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.font.size.pt == pytest.approx(9.75)


# ---------------------------------------------------------------------------
# Legend position (Fix 5)
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "position,expected",
    [
        ("right", XL_LEGEND_POSITION.RIGHT),
        ("left", XL_LEGEND_POSITION.LEFT),
        ("top", XL_LEGEND_POSITION.TOP),
        ("bottom", XL_LEGEND_POSITION.BOTTOM),
    ],
)
def test_legend_position_mapped(position, expected):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "legend": {"show": True, "position": position},
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.legend.position == expected


@pytest.mark.parametrize("position", [None, "corner", "diagonal", 42])
def test_legend_position_defaults_to_bottom_when_absent_or_unrecognized(position):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "legend": {"show": True, "position": position},
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.legend.position == XL_LEGEND_POSITION.BOTTOM


# ---------------------------------------------------------------------------
# Fill transparency / alpha (Fix 4)
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "value,expected",
    [
        ("rgba(53, 208, 186, 0.28)", 0.28),
        ("rgba(0,0,0,0)", 0.0),
        ("rgba(1,2,3,1)", None),
        ("rgba(1,2,3,1.0)", None),
        ("rgb(1,2,3)", None),
        ("#FF0000", None),
        ("rgba(1,2,3,bad)", None),
        ("rgba(1,2,3,2)", None),
        ("rgba(1,2,3,-1)", None),
        (None, None),
        (123, None),
    ],
)
def test_alpha_from_css_color(value, expected):
    result = svc._alpha_from_css_color(value)
    if expected is None:
        assert result is None
    else:
        assert result == pytest.approx(expected)


@pytest.mark.parametrize(
    "text,expected_hex",
    [
        ("#fff", "FFFFFF"),
        ("#123456", "123456"),
        ("rgb(10, 20, 30)", "0A141E"),
        ("rgba(10, 20, 30, 0.5)", "0A141E"),
        ("not-a-color", None),
        (None, None),
    ],
)
def test_hex_from_css_color_unchanged(text, expected_hex):
    """Pins _hex_from_css_color's existing contract - _alpha_from_css_color
    was added as a separate function specifically so this signature/behavior
    never had to change."""
    assert svc._hex_from_css_color(text) == expected_hex


def test_area_series_fill_preserves_rgba_alpha():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "area",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B", "C"],
        "datasets": [
            {
                "label": "S1",
                "data": [1, 2, 3],
                "backgroundColor": "rgba(53, 208, 186, 0.28)",
            }
        ],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    series = chart.plots[0].series[0]
    assert series.format.fill.fore_color.rgb == RGBColor.from_string("35D0BA")
    srgbClr = series.format.fill.fore_color._xFill.find(qn("a:srgbClr"))
    alpha = srgbClr.find(qn("a:alpha"))
    assert alpha is not None
    assert alpha.get("val") == "28000"


def test_per_point_fill_preserves_rgba_alpha_independently():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [
            {
                "label": "S1",
                "data": [1, 2],
                "backgroundColor": ["rgba(255,0,0,0.5)", "#00FF00"],
            }
        ],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    points = list(chart.plots[0].series[0].points)

    p0_srgbClr = points[0].format.fill.fore_color._xFill.find(qn("a:srgbClr"))
    p0_alpha = p0_srgbClr.find(qn("a:alpha"))
    assert p0_alpha is not None
    assert p0_alpha.get("val") == "50000"

    p1_srgbClr = points[1].format.fill.fore_color._xFill.find(qn("a:srgbClr"))
    assert p1_srgbClr.find(qn("a:alpha")) is None


def test_opaque_colors_emit_no_alpha():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "area",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [
            {"label": "S1", "data": [1, 2], "backgroundColor": "rgba(51, 102, 255, 1)"}
        ],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart._chartSpace.find(".//" + qn("a:alpha")) is None


def test_alpha_survives_save_and_reload(tmp_path):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "area",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [
            {
                "label": "S1",
                "data": [1, 2],
                "backgroundColor": "rgba(53, 208, 186, 0.28)",
            }
        ],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    out_path = tmp_path / "area.pptx"
    prs.save(str(out_path))
    reopened = Presentation(str(out_path))
    chart2 = reopened.slides[0].shapes[0].chart
    srgbClr = chart2.plots[0].series[0].format.fill.fore_color._xFill.find(qn("a:srgbClr"))
    alpha = srgbClr.find(qn("a:alpha"))
    assert alpha is not None
    assert alpha.get("val") == "28000"


def test_chart_and_plot_area_are_made_transparent():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart

    chartSpace = chart._chartSpace
    chart_el = chartSpace.find(qn("c:chart"))
    plotArea = chart_el.find(qn("c:plotArea"))
    # The chart-area spPr is a sibling of c:chart (a direct child of
    # c:chartSpace, inserted right after c:chart), not a child of c:chart.
    assert chart_el.getnext() is not None
    assert chart_el.getnext().tag == qn("c:spPr")
    assert plotArea.find(qn("c:spPr")) is not None
    assert len(chartSpace.findall(".//" + qn("a:noFill"))) >= 2


def test_transparency_survives_save_and_reload():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "pie",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2], "backgroundColor": ["#111111", "#222222"]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    import tempfile, os

    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "out.pptx")
        prs.save(path)
        reopened = Presentation(path)
        chart2 = reopened.slides[0].shapes[0].chart
        chartSpace2 = chart2._chartSpace
        chart_el2 = chartSpace2.find(qn("c:chart"))
        assert chart_el2.find(qn("c:plotArea")).find(qn("c:spPr")) is not None
        assert chart_el2.getnext() is not None and chart_el2.getnext().tag == qn("c:spPr")


def _donut_capture(**overrides):
    captured = {
        "kind": "donut",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["OpenAI", "Google", "Anthropic"],
        "datasets": [
            {
                "label": "S1",
                "data": [38, 22, 15],
                "backgroundColor": ["#E00600", "#0B1F3A", "#000000"],
            }
        ],
        "hasDataLabels": True,
    }
    captured.update(overrides)
    return captured


def test_pie_data_labels_show_category_name_and_percentage_outside_end():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = _donut_capture(kind="pie")
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    data_labels = chart.plots[0].data_labels
    assert data_labels.show_category_name is True
    assert data_labels.show_percentage is True
    assert data_labels.show_value is False
    # Legal on pie - must be present so labels fit outside slim slices.
    dLblPos = chart._chartSpace.find(".//" + qn("c:dLblPos"))
    assert dLblPos is not None
    assert dLblPos.get("val") == "outEnd"


def test_pie_data_labels_ignore_the_browser_captured_color():
    """Real reported bug: a pie chart's labels were present in the exported
    PPTX but invisible - white text on the white page background. Cause:
    pie labels are repositioned OUTSIDE the slice on export (see the test
    above - PowerPoint requires this to legally fit "Category NN%" outside a
    slim slice), but the captured `dataLabelColor` was chosen in the browser
    for a label sitting INSIDE a colored slice (this app's own prompt asks
    for white text there specifically), so applying it verbatim to the
    now-outside, now-against-white-background label produced white-on-white.
    Pie must leave the color unset - same "let PowerPoint pick a sane
    default" approach already used for donut's data-label position."""
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = _donut_capture(kind="pie", dataLabelColor="#FFFFFF")
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    data_labels = chart.plots[0].data_labels
    with pytest.raises(AttributeError):
        # python-pptx raises when no color has ever been set - this is the
        # "unset, PowerPoint picks its own default" state, not a color that
        # merely happens to be black/dark.
        _ = data_labels.font.color.rgb


def test_donut_data_labels_still_get_the_captured_color():
    """Contrast case for the fix above: donut labels stay INSIDE the ring on
    export (no OUTSIDE_END repositioning), so the browser-captured
    inside-slice color is still correct there and must still be applied -
    the pie fix must not accidentally suppress this for every circular
    chart kind."""
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = _donut_capture(dataLabelColor="#FFFFFF")
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    data_labels = chart.plots[0].data_labels
    assert data_labels.font.color.rgb == RGBColor.from_string("FFFFFF")


def test_donut_data_labels_show_percentage_only_and_never_emit_dLblPos():
    """PowerPoint rejects <c:dLblPos> on a doughnut chart outright - a real
    export with it set failed to open and PowerPoint's repair silently
    dropped the entire slide. Must assert on the emitted XML, not just the
    python-pptx property, since python-pptx itself performs no validation
    and would happily accept the value that corrupts the file."""
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = _donut_capture()
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    data_labels = chart.plots[0].data_labels
    assert data_labels.show_category_name is False
    assert data_labels.show_percentage is True
    assert data_labels.show_value is False
    assert chart._chartSpace.find(".//" + qn("c:dLblPos")) is None


def test_donut_dLblPos_stays_absent_after_save_and_reload(tmp_path):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = _donut_capture()
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    out_path = tmp_path / "donut.pptx"
    prs.save(str(out_path))
    reopened = Presentation(str(out_path))
    chart = list(reopened.slides[0].shapes)[0].chart
    assert chart._chartSpace.find(".//" + qn("c:dLblPos")) is None


@pytest.mark.parametrize(
    "kind",
    [
        "bar",
        "horizontal_bar",
        "stacked_bar",
        "horizontal_stacked_bar",
        "line",
        "area",
        "radar",
        "donut",
    ],
)
def test_no_kind_except_pie_emits_dLblPos(kind):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": kind,
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B", "C"],
        "datasets": [{"label": "S1", "data": [1, 2, 3]}],
        "hasDataLabels": True,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart._chartSpace.find(".//" + qn("c:dLblPos")) is None


def test_bar_data_labels_show_value_not_category_or_percentage():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "hasDataLabels": True,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    data_labels = chart.plots[0].data_labels
    assert data_labels.show_category_name is False
    assert data_labels.show_percentage is False
    # Regression: plot.has_data_labels = True alone emits <c:showVal val="0"/>
    # for every non-circular kind, so a real bar/line/area/radar chart's data
    # labels displayed nothing at all until this was explicitly turned on.
    assert data_labels.show_value is True
    showVal = chart._chartSpace.find(".//" + qn("c:showVal"))
    assert showVal is not None
    assert showVal.get("val") == "1"


@pytest.mark.parametrize("kind", ["bar", "line", "area"])
def test_non_circular_kinds_show_value_true(kind):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": kind,
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B", "C"],
        "datasets": [{"label": "S1", "data": [1, 2, 3]}],
        "hasDataLabels": True,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    data_labels = list(slide.shapes)[0].chart.plots[0].data_labels
    assert data_labels.show_value is True


def test_area_data_labels_never_get_a_position_real_corruption_found():
    """Regression test for a real, user-confirmed corruption: an earlier
    version of this code set `data_labels.position = XL_LABEL_POSITION.CENTER`
    for area charts specifically, reasoning that CENTER ("ctr") is the one
    position value Microsoft's own documentation claims is valid across every
    chart type with no exception (unlike OUTSIDE_END/etc, bar/pie only, or
    ABOVE/BELOW/LEFT/RIGHT, line/scatter/radar only). A user opened the real
    exported file in real PowerPoint and got a "needs repair" prompt - repair
    deleted the area chart's own slide, the same failure mode as the
    dLblPos-on-doughnut bug this was meant to avoid repeating. So "ctr" is
    apparently *also* illegal for area in real PowerPoint, contradicting the
    documentation - python-pptx itself raised no error and re-read the file
    fine both times, so only a real PowerPoint open ever caught this. Lesson:
    an area chart's data labels must never get any explicit position at all;
    this is now a hard invariant, not a "which value is safe" question."""
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "area",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B", "C"],
        "datasets": [{"label": "S1", "data": [1, 2, 3]}],
        "hasDataLabels": True,
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    assert chart.plots[0].data_labels.position is None
    assert chart._chartSpace.find(".//" + qn("c:dLblPos")) is None


def test_axis_title_color_applied():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 400, 300)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 400, "height": 300},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
        "axisTitles": {"x": "Category", "y": "Value"},
        "axisTitleColor": "#AFC3CE",
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True
    chart = list(slide.shapes)[0].chart
    cat_color = chart.category_axis.axis_title.text_frame.paragraphs[0].font.color.rgb
    val_color = chart.value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb
    assert cat_color == RGBColor.from_string("AFC3CE")
    assert val_color == RGBColor.from_string("AFC3CE")


def test_try_upgrade_one_chart_survives_save_and_reload():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 500, 400)
    captured_chart = {
        "kind": "pie",
        "hasMarkers": False,
        "boundingBox": {"left": 0, "top": 0, "width": 500, "height": 400},
        "labels": ["A", "B"],
        "datasets": [{"label": "S1", "data": [1, 2], "backgroundColor": ["#FF0000", "#00FF00"]}],
        "title": None,
        "legend": {"show": False},
        "axisTitles": {},
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    import tempfile, os

    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "out.pptx")
        prs.save(path)
        reopened = Presentation(path)
        shapes = list(reopened.slides[0].shapes)
        assert len(shapes) == 1
        assert shapes[0].has_chart
        assert shapes[0].chart.chart_type == XL_CHART_TYPE.PIE


def test_try_upgrade_one_chart_unsupported_kind_leaves_picture_untouched():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 100, 100)
    captured_chart = {
        "kind": "scatter",
        "boundingBox": {"left": 0, "top": 0, "width": 100, "height": 100},
        "labels": ["a"],
        "datasets": [{"label": "S1", "data": [1]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is False
    shapes = list(slide.shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_try_upgrade_one_chart_no_confident_match_leaves_picture_untouched():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 100, 100)
    # Bounding box far away from the actual picture -> no overlap -> no match.
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 900, "top": 600, "width": 100, "height": 100},
        "labels": ["a"],
        "datasets": [{"label": "S1", "data": [1]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is False
    shapes = list(slide.shapes)
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_try_upgrade_one_chart_malformed_data_leaves_picture_untouched():
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(0, 0, 100, 100)
    captured_chart = {
        "kind": "bar",
        "boundingBox": {"left": 0, "top": 0, "width": 100, "height": 100},
        "labels": ["a", "b"],
        "datasets": [{"label": "S1", "data": [1]}],  # length mismatch -> _build_chart_data -> None
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is False
    shapes = list(slide.shapes)
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_try_upgrade_one_chart_preserves_z_order():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    emu_x = prs.slide_width / 1280
    emu_y = prs.slide_height / 720

    before = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(100), Emu(100))
    before.text_frame.text = "before"

    img = Image.new("RGB", (10, 10), color="blue")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(
        buf, Emu(int(100 * emu_x)), Emu(int(100 * emu_y)), Emu(int(200 * emu_x)), Emu(int(200 * emu_y))
    )

    after = slide.shapes.add_textbox(Emu(0), Emu(300), Emu(100), Emu(100))
    after.text_frame.text = "after"

    captured_chart = {
        "kind": "line",
        "hasMarkers": True,
        "boundingBox": {"left": 100, "top": 100, "width": 200, "height": 200},
        "labels": ["a", "b"],
        "datasets": [{"label": "S1", "data": [1, 2]}],
    }
    claimed: set = set()
    assert svc._try_upgrade_one_chart(slide, captured_chart, emu_x, emu_y, claimed) is True

    shape_types = [s.shape_type for s in slide.shapes]
    assert shape_types == [MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TEXT_BOX]


# ---------------------------------------------------------------------------
# Top-level orchestration / non-regression safety net
# ---------------------------------------------------------------------------


def test_upgrade_flattened_charts_to_native_noop_without_token():
    # Must never raise, and must not touch the file, when no token is given.
    asyncio.run(
        svc.upgrade_flattened_charts_to_native("/nonexistent/path.pptx", None, uuid.uuid4())
    )


def test_upgrade_flattened_charts_to_native_noop_when_nothing_captured(monkeypatch):
    monkeypatch.setattr(svc.chart_capture_store, "take_capture", lambda token: None)
    # Should return immediately without ever touching the (nonexistent) file.
    asyncio.run(
        svc.upgrade_flattened_charts_to_native(
            "/nonexistent/path.pptx", "some-token", uuid.uuid4()
        )
    )


def test_upgrade_flattened_charts_to_native_end_to_end(monkeypatch, tmp_path):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(100, 100, 400, 300)
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))

    captured = {
        "presentation_id": "irrelevant",
        "charts": [
            {
                "slideOrderIndex": 0,
                "kind": "bar",
                "hasMarkers": False,
                "boundingBox": {"left": 100, "top": 100, "width": 400, "height": 300},
                "labels": ["Q1", "Q2"],
                "datasets": [{"label": "S1", "data": [1, 2]}],
            }
        ],
    }
    monkeypatch.setattr(
        svc.chart_capture_store, "take_capture", lambda token: captured
    )

    async def _fake_expected_count(presentation_id):
        return 1

    monkeypatch.setattr(svc, "_expected_slide_count", _fake_expected_count)

    asyncio.run(
        svc.upgrade_flattened_charts_to_native(str(pptx_path), "tok", uuid.uuid4())
    )

    reopened = Presentation(str(pptx_path))
    shapes = list(reopened.slides[0].shapes)
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.CHART


def test_upgrade_flattened_charts_to_native_slide_count_mismatch_is_untouched(
    monkeypatch, tmp_path
):
    prs, slide, picture, emu_x, emu_y = _make_presentation_with_picture(100, 100, 400, 300)
    pptx_path = tmp_path / "deck.pptx"
    prs.save(str(pptx_path))
    original_bytes = pptx_path.read_bytes()

    captured = {
        "presentation_id": "irrelevant",
        "charts": [
            {
                "slideOrderIndex": 0,
                "kind": "bar",
                "boundingBox": {"left": 100, "top": 100, "width": 400, "height": 300},
                "labels": ["Q1"],
                "datasets": [{"label": "S1", "data": [1]}],
            }
        ],
    }
    monkeypatch.setattr(svc.chart_capture_store, "take_capture", lambda token: captured)

    async def _fake_expected_count(presentation_id):
        return 2  # deliberately wrong -> should bail out entirely

    monkeypatch.setattr(svc, "_expected_slide_count", _fake_expected_count)

    asyncio.run(svc.upgrade_flattened_charts_to_native(str(pptx_path), "tok", uuid.uuid4()))

    assert pptx_path.read_bytes() == original_bytes
